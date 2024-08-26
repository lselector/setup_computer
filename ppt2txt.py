
"""
# ppt2txt.py
# script goes recursively through directory and subdirectories
# and extracts text from all PowerPoint pptx files
# and saves to a text file 
# It looks either in current directory
# or in directories provided on command line (space-separaed)
#
# pip install python-pptx
#
# Example of usage:
# in ml_ai_doc directory:
#     python ppt2txt.py
# in other directory:
#     python  ../setup_computer/ppt2txt.py ./
# you can put this command into extract.bash in each directory
""" 
# --------------------------------------------------------------
import os, sys, glob
import datetime as dt
from pathlib import Path
from pptx import Presentation

# --------------------------------------------------------------
def print_date_time():
    print (dt.datetime.now().strftime("%Y-%m-%d %H:%M:%S"))

print_date_time()

# --------------------------------------------------------------
if len(sys.argv) <= 1:
    mypath = os.path.dirname(os.path.realpath(__file__))
    print("looking for pptx files in script directory: ", mypath)
    mydirs = [mypath]
else:
    mydirs = []
    for argpath in sys.argv[1:]:
        if not os.path.isdir(argpath):
            print(f"{argpath} is not a directory, exiting ...")
            sys.exit()
        mypath = os.path.realpath(argpath)
        print(argpath, "   =>  ", mypath)
        mydirs += [mypath]

# --------------------------------------------------------------
for dir_path in mydirs:
    print_date_time()
    print(dir_path)
    basedir = dir_path.split("/")[-1]
    fname_out = dir_path + "/" + basedir+"_ppt.txt"

    # get paths recursively 
    myfiles = list(Path(dir_path).rglob("*.[pP][pP][tT][xX]"))

    # convert from posix path to string
    myfiles = [str(x) for x in myfiles] 

    counter = 0
    ss = ""

    for myfile in sorted(myfiles):
        if "~$" in myfile:
            continue
        if "OLD/" in myfile:
            continue
        counter += 1
        if counter > 3:
            pass # break
        print_date_time()
        prs = Presentation(myfile)
        print("-"*20, myfile)
        myfile_short = myfile.split(dir_path)[1][1:]
        # print(myfile_short,"\n")
        # print("----------------------")
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    txt = shape.text
                    txt_arr = txt.strip().split("\n")
                    txt_arr = [x.strip() for x in txt_arr]
                    txt_arr = [x.replace("\r"," ") for x in txt_arr]
                    txt_arr = [myfile_short+" : " + x + "\n" for x in txt_arr]
                    txt = "".join(txt_arr)
                    ss += txt
                    print(txt)
    print("----------------------")
    print("writing to file ", fname_out)
    fh = open(fname_out, "w")
    fh.write(ss)
    fh.close()
    print_date_time()
    print("----------------------")
